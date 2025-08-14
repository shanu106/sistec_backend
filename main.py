import os
import json
import re
from datetime import datetime
from io import BytesIO
from storage3.utils import StorageException 
import requests
import httpx
from fastapi import FastAPI, Form, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pollinations

from supabase import create_client, Client
os.environ["HTTPX_HTTP2"] = "false"

# ===== CONFIG =====
GEMINI_API_KEY = "AIzaSyAAd7Q2eyvgaAc3LLkKdXHHBpt2wpUaLWM"
UNSPLASH_ACCESS_KEY = "NFqWtuuBhhyq84vDmMfzpgNd1PvTuahe8-bXrBxHWUw"
GEMINI_MODEL_TEXT = "gemini-2.0-flash"
GEMINI_ENDPOINT = "https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"

# Supabase Config - set your env variables SUPABASE_URL and SUPABASE_KEY accordingly
SUPABASE_URL = "https://dqndyittfzlxufcysuma.supabase.co"
# SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6ImRxbmR5aXR0ZnpseHVmY3lzdW1hIiwicm9sZSI6ImFub24iLCJpYXQiOjE3NTUwNDMzMzMsImV4cCI6MjA3MDYxOTMzM30.kvvqKnohH9LOE4mgeSk8Gx25Cou1u1zEReWFLo2CSzA"
SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6ImRxbmR5aXR0ZnpseHVmY3lzdW1hIiwicm9sZSI6InNlcnZpY2Vfcm9sZSIsImlhdCI6MTc1NTA0MzMzMywiZXhwIjoyMDcwNjE5MzMzfQ.YY0TbVzkYv5h_mbgfYCmEkLY6ypARgEoB9jqcS211jA"
BUCKET_NAME = "ppt-files"

# Create an HTTPX client with HTTP/2 disabled to avoid SSLV3_ALERT_BAD_RECORD_MAC errors
client = httpx.Client(http2=False, timeout=30.0)
supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

app = FastAPI(title="AI PPT Generator with Supabase Storage", version="1.0")

# ===== DEFAULT STYLE =====
DEFAULT_STYLE = {
    "bg_color": RGBColor(9, 16, 96),
    "title_color": RGBColor(13, 238, 231),
    "text_color": RGBColor(255, 255, 255),
    "font": "Calibri"
}

# ===== IMAGE GENERATION (returns bytes) =====
def generate_image_with_pollinations(prompt: str) -> bytes:
    model = pollinations.Image(
        model="flux",
        width=1024,
        height=1024,
        seed=42
    )
    img = model.Generate(prompt=prompt, save=False)
    img_byte_arr = BytesIO()
    img.save(img_byte_arr, format='JPEG')
    img_byte_arr.seek(0)
    return img_byte_arr.read()

# ===== FALLBACK IMAGE (returns bytes or None) =====
def get_fallback_image(query) -> bytes | None:
    url = f"https://api.unsplash.com/photos/random?query={query}&orientation=landscape&client_id={UNSPLASH_ACCESS_KEY}"
    resp = requests.get(url)
    if resp.status_code == 200:
        img_url = resp.json()["urls"]["regular"]
        img_data = requests.get(img_url).content
        return img_data
    return None

# ===== STYLE EXTRACTION =====
def extract_style_from_prompt(prompt):
    ai_prompt = f"""
From the text: "{prompt}", detect PPT style.
Respond ONLY in JSON:
{{
  "bg_color": "hex",
  "title_color": "hex",
  "text_color": "hex",
  "font": "string"
}}
If not provided, return defaults.
"""
    url = GEMINI_ENDPOINT.format(model=GEMINI_MODEL_TEXT)
    resp = requests.post(
        url,
        headers={"Content-Type": "application/json"},
        params={"key": GEMINI_API_KEY},
        json={"contents": [{"parts": [{"text": ai_prompt}]}]},
        timeout=30
    )
    if resp.status_code != 200:
        return DEFAULT_STYLE

    txt = resp.json().get("candidates", [{}])[0].get("content", {}).get("parts", [{"text": ""}])[0]["text"].strip()
    match = re.search(r"\{.*\}", txt, re.S)
    if match:
        txt = match.group(0)

    try:
        style_data = json.loads(txt)
        return {
            "bg_color": RGBColor(*(int(style_data["bg_color"].lstrip("#")[i:i + 2], 16) for i in (0, 2, 4))),
            "title_color": RGBColor(*(int(style_data["title_color"].lstrip("#")[i:i + 2], 16) for i in (0, 2, 4))),
            "text_color": RGBColor(*(int(style_data["text_color"].lstrip("#")[i:i + 2], 16) for i in (0, 2, 4))),
            "font": style_data.get("font", "Calibri")
        }
    except:
        return DEFAULT_STYLE

# ===== GEMINI SLIDE GENERATOR =====
def get_slides_from_gemini(topic):
    prompt = f"""
You are a presentation designer. Output ONLY valid JSON with:
{{
  "slides":[
    {{
      "title":"string",
      "content":"string (newline bullets)",
      "image_prompt":"string (detailed prompt for image generation)"
    }}
  ]
}}
Topic: {topic}
Only JSON, no extra text.
"""
    url = GEMINI_ENDPOINT.format(model=GEMINI_MODEL_TEXT)
    resp = requests.post(
        url,
        headers={"Content-Type": "application/json"},
        params={"key": GEMINI_API_KEY},
        json={"contents": [{"parts": [{"text": prompt}]}]},
        timeout=30
    )
    txt = resp.json().get("candidates", [{}])[0].get("content", {}).get("parts", [{"text": ""}])[0]["text"]

    txt = txt.strip().strip("```").lstrip("json").strip()
    return json.loads(txt)

# ===== PPT BUILDER (returns PPTX bytes) =====
def build_presentation(data, style) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Cover slide
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    fill = cover.background.fill
    fill.solid()
    fill.fore_color.rgb = style["bg_color"]

    title_box = cover.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = data.get("slides", [])[0].get("title", "Presentation")
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.name = style["font"]
    p.font.color.rgb = style["title_color"]
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = cover.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1))
    tf_sub = subtitle_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.text = "Generated by AI Alliance"
    p_sub = tf_sub.paragraphs[0]
    p_sub.font.size = Pt(20)
    p_sub.font.bold = False
    p_sub.font.name = style["font"]
    p_sub.font.color.rgb = style["text_color"]
    p_sub.alignment = PP_ALIGN.CENTER

    # Content slides
    for slide in data.get("slides", []):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = style["bg_color"]

        title_box = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(6), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = slide.get("title", "")
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = style["font"]
        p.font.color.rgb = style["title_color"]

        content_box = s.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(6), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(slide.get("content", "").splitlines()):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line.strip()
            p.font.size = Pt(18)
            p.font.name = style["font"]
            p.font.color.rgb = style["text_color"]

        img_bytes = None
        try:
            img_bytes = generate_image_with_pollinations(slide.get("image_prompt", ""))
        except Exception as e:
            print(f"[WARN] Pollinations image fail for '{slide['title']}': {e}")

        if not img_bytes:
            img_bytes = get_fallback_image(slide.get("image_prompt", slide['title']))

        if img_bytes:
            img_stream = BytesIO(img_bytes)
            s.shapes.add_picture(img_stream, Inches(7.2), Inches(1.5), width=Inches(5.5))

    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream.read()

# ===== Upload PPT to Supabase Storage =====
def upload_ppt_to_supabase(filename: str, data: bytes) -> str:
    try:
        filename = filename.strip().replace("\n", "")
        supabase.storage.from_(BUCKET_NAME).upload(filename, data)
    except Exception as e:
        raise Exception(f"Upload failed: {e}")

    # Build public URL (works if bucket is public)
    file_url = f"{SUPABASE_URL}/storage/v1/object/public/{BUCKET_NAME}/{filename}"
    return file_url

 
 

#
 
 
    
    
    
@app.post("/generate_ppt")
async def generate_ppt_api(topic: str = Form(...)):
    print("function started :")
    style = extract_style_from_prompt(topic)
    print("style got")
    slides_data = get_slides_from_gemini(topic)
    print("slides data got")
    pptx_bytes = build_presentation(slides_data, style)
    print("presentation built")
    safe_filename = re.sub(r'[\\/*?:"<>|]', "", topic).replace(" ", "_").strip()
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M")
    filename = f"{safe_filename}_{timestamp}.pptx"
    print(f"Generated filename: {filename}")
    file_url = upload_ppt_to_supabase(filename, pptx_bytes)
    print("file uploaded to Supabase")
    if not file_url:
        raise HTTPException(status_code=500, detail="Failed to upload file to Supabase Storage")

    return {"message": "PPT generated and uploaded", "filename": filename, "file_url": file_url}

@app.get("/download_ppt/{filename}")
async def download_ppt(filename: str):
    response = supabase.storage.from_(BUCKET_NAME).download(filename)
    if response.get("error") or not response.get("data"):
        raise HTTPException(status_code=404, detail="File not found")
    file_bytes = response["data"]
    return StreamingResponse(BytesIO(file_bytes),
                             media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                             headers={"Content-Disposition": f"attachment; filename={filename}"})